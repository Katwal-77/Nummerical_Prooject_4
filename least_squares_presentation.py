from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a new presentation
prs = Presentation()

# Set slide dimensions to 16:9 aspect ratio
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Title slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Discovering Laws Behind Real-World Data Using the Least Squares Method"
subtitle.text = "Prem Bahadur Katuwal\nStudent ID: 202424080129\nNumerical Analysis (PhD Level)"

# Introduction slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Introduction"
content.text = "• The least squares method is a fundamental technique for discovering patterns in data\n\n• Developed by Gauss and Legendre in the early 19th century\n\n• Minimizes the sum of squared differences between observed and predicted values\n\n• Applications: physics, engineering, economics, social sciences, etc.\n\n• This presentation demonstrates its application to real-world data"

# Theoretical Foundation slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Theoretical Foundation"
content.text = "• Mathematical Formulation:\n   S(β) = Σ[yi - f(xi; β)]²\n\n• Matrix Approach for Linear Models:\n   S(β) = ||y - Xβ||²\n\n• Normal Equations:\n   X^T X β = X^T y\n\n• Solution:\n   β = (X^T X)^(-1) X^T y"

# Linear Least Squares slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Linear Least Squares: House Size vs. Price"
content.text = "• Dataset: House sizes (sqft) and prices ($1000s)\n\n• Linear model: Price = m × Size + c\n\n• Best-fit parameters:\n   - Slope (m): 0.0991\n   - Intercept (c): 1.9382\n\n• Discovered Law:\n   Price = 0.0991 × Size + 1.9382 ($1000s)\n\n• R² = 0.9980 (99.8% of variance explained)"

# Linear Least Squares - Visualization slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Linear Least Squares: Visualization"

# Add images
left = Inches(1)
top = Inches(1.5)
width = Inches(5)
height = Inches(3.75)
slide.shapes.add_picture("fitted_line.png", left, top, width, height)

left = Inches(7)
slide.shapes.add_picture("residuals.png", left, top, width, height)

# Add captions
left = Inches(2)
top = Inches(5.5)
width = Inches(3)
height = Inches(0.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Fitted Line (R² = 0.9980)"
p.alignment = PP_ALIGN.CENTER

left = Inches(8)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Residuals Analysis"
p.alignment = PP_ALIGN.CENTER

# Polynomial Least Squares slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Polynomial Least Squares: Temperature Data"
content.text = "• Dataset: Monthly temperatures in northern hemisphere\n\n• Polynomial model: T = β₀ + β₁x + β₂x² + ... + βₚxᵖ\n\n• Best-fit (4th degree):\n   T = 0.021x⁴ - 0.595x³ + 4.900x² - 9.137x - 0.417\n\n• R² = 0.9984 (99.84% of variance explained)\n\n• Captures seasonal temperature variations"

# Polynomial Least Squares - Visualization slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Polynomial Least Squares: Visualization"

# Add images
left = Inches(1)
top = Inches(1.5)
width = Inches(5)
height = Inches(3.75)
slide.shapes.add_picture("temperature_fit.png", left, top, width, height)

left = Inches(7)
slide.shapes.add_picture("temperature_r2.png", left, top, width, height)

# Add captions
left = Inches(2)
top = Inches(5.5)
width = Inches(3)
height = Inches(0.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "4th Degree Polynomial Fit"
p.alignment = PP_ALIGN.CENTER

left = Inches(8)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Model Selection: R² vs. Degree"
p.alignment = PP_ALIGN.CENTER

# Non-Linear Least Squares slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Non-Linear Least Squares: Population Growth"
content.text = "• Dataset: World population from 1900 to 2020\n\n• Exponential model: P(t) = P₀e^(rt)\n\n• Two approaches:\n   1. Linearization: ln(P) = ln(P₀) + rt\n   2. Direct non-linear fitting\n\n• Discovered Law:\n   P(t) = 1.69e^(0.013t), where t = years since 1900\n\n• Growth rate: ~1.3% per year (doubling time: ~54 years)"

# Non-Linear Least Squares - Visualization slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Non-Linear Least Squares: Visualization"

# Add image
left = Inches(4)
top = Inches(1.5)
width = Inches(5)
height = Inches(3.75)
slide.shapes.add_picture("population_growth.png", left, top, width, height)

# Add caption
left = Inches(5)
top = Inches(5.5)
width = Inches(3)
height = Inches(0.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Exponential Growth Models"
p.alignment = PP_ALIGN.CENTER

# Advanced Analysis slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Advanced Analysis"
content.text = "• Residual Analysis:\n   - Random scatter around zero\n   - No obvious patterns\n   - No extreme outliers\n\n• Confidence Intervals:\n   - Quantify uncertainty in predictions\n   - Wider at extremes of data range\n\n• Goodness of Fit Metrics:\n   - R² (coefficient of determination)\n   - RMSE (root mean square error)\n   - MAE (mean absolute error)"

# Confidence Intervals slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Confidence Intervals"

# Add image
left = Inches(4)
top = Inches(1.5)
width = Inches(5)
height = Inches(3.75)
slide.shapes.add_picture("confidence_intervals.png", left, top, width, height)

# Add caption
left = Inches(5)
top = Inches(5.5)
width = Inches(3)
height = Inches(0.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "95% Confidence Intervals for Predictions"
p.alignment = PP_ALIGN.CENTER

# Comparison with Other Methods slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Comparison with Other Methods"
content.text = "• Robust Regression:\n   - Less sensitive to outliers\n   - Uses Huber loss or Tukey's bisquare\n\n• Ridge Regression:\n   - Adds penalty term to handle multicollinearity\n   - Minimizes Σ[yi - f(xi; β)]² + λΣβj²\n\n• LASSO Regression:\n   - Encourages sparse solutions\n   - Minimizes Σ[yi - f(xi; β)]² + λΣ|βj|\n\n• Orthogonal Distance Regression:\n   - Accounts for errors in both variables"

# Conclusion slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Conclusion"
content.text = "• The least squares method is powerful for discovering laws in data\n\n• Key findings from our analysis:\n   - House price follows a linear law: Price ≈ 0.10 × Size + 1.94 ($1000s)\n   - Monthly temperatures follow a 4th-degree polynomial pattern\n   - Population growth follows an exponential law: P(t) = 1.69e^(0.013t)\n\n• Applications:\n   - Prediction and forecasting\n   - Understanding relationships\n   - Testing hypotheses\n   - Decision-making"

# Thank You slide
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Thank You!"

# Add text box
left = Inches(2)
top = Inches(3)
width = Inches(9.33)
height = Inches(1.5)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Prem Bahadur Katuwal\nStudent ID: 202424080129\nNumerical Analysis (PhD Level)"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(32)
p.font.bold = True

# Save the presentation
prs.save('least_squares_presentation.pptx')
